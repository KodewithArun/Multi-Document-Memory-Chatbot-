import streamlit as st
import os
from pathlib import Path
import pandas as pd
from pptx import Presentation
import json
from datetime import datetime
import time
import gc
import shutil
import tempfile

# Import all the necessary libraries from your original code
from llama_parse import LlamaParse
from dotenv import load_dotenv
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_groq import ChatGroq
from langchain_community.document_loaders import (
    UnstructuredPDFLoader,
    UnstructuredWordDocumentLoader,
    UnstructuredPowerPointLoader,
    UnstructuredExcelLoader,
    UnstructuredMarkdownLoader,
)
from langchain_community.document_loaders import TextLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from langchain.retrievers import SelfQueryRetriever
from langchain.chains.query_constructor.base import AttributeInfo
from langchain.prompts.chat import (
    ChatPromptTemplate,
    MessagesPlaceholder,
    SystemMessagePromptTemplate,
    HumanMessagePromptTemplate,
)
from langchain.schema.runnable import RunnableSequence, RunnableLambda
from langchain_core.runnables.history import RunnableWithMessageHistory
from langchain_community.chat_message_histories import FileChatMessageHistory
from langchain_core.output_parsers import StrOutputParser

# Load environment variables
load_dotenv()

# Initialize session state
if 'vectorstore' not in st.session_state:
    st.session_state.vectorstore = None
if 'retrieval_chain_with_memory' not in st.session_state:
    st.session_state.retrieval_chain_with_memory = None
if 'session_id' not in st.session_state:
    st.session_state.session_id = f"user_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []

# Constants
SUPPORTED_EXTS = [".pdf", ".docx", ".pptx", ".md", ".txt", ".xlsx"]

# Utility functions for safe file operations
def safe_delete_file(file_path):
    """Safely delete a file with retry logic and proper error handling"""
    if not file_path or not Path(file_path).exists():
        return True
    
    max_attempts = 5
    for attempt in range(max_attempts):
        try:
            # Force garbage collection to release file handles
            gc.collect()
            time.sleep(0.1)  # Small delay to ensure file handles are released
            
            # Convert to Path object if it's a string
            path_obj = Path(file_path)
            if path_obj.exists():
                path_obj.unlink()
            return True
            
        except PermissionError:
            if attempt < max_attempts - 1:
                time.sleep(1)  # Wait longer between retries
                continue
            else:
                st.warning(f"Could not delete {file_path} after {max_attempts} attempts. File may be in use.")
                return False
        except Exception as e:
            st.error(f"Unexpected error deleting {file_path}: {e}")
            return False
    return True

def safe_read_excel(file_path):
    """Safely read Excel file with proper error handling"""
    try:
        # Use context manager to ensure file is properly closed
        with pd.ExcelFile(str(file_path)) as xls:
            content = ""
            for sheet_name in xls.sheet_names:
                try:
                    df = pd.read_excel(xls, sheet_name=sheet_name)
                    if df.empty:
                        continue
                    content += f"## Sheet: {sheet_name}\n\n"
                    content += df.to_markdown(index=False) + "\n\n"
                except Exception as e:
                    st.warning(f"Error reading sheet {sheet_name}: {e}")
                    continue
            return content.strip() if content.strip() else None
    except Exception as e:
        st.error(f"XLSX read error: {e}")
        return None

def safe_read_file(file_path, encoding="utf-8"):
    """Safely read text file with proper error handling"""
    try:
        with open(file_path, 'r', encoding=encoding) as f:
            content = f.read()
        return content.strip() if content.strip() else None
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='latin-1') as f:
                content = f.read()
            return content.strip() if content.strip() else None
        except Exception as e:
            st.error(f"Error reading file {file_path}: {e}")
            return None
    except Exception as e:
        st.error(f"Error reading file {file_path}: {e}")
        return None

# User-specific folders
def get_user_folders(session_id):
    base_folder = Path("user_data") / session_id
    parsed_folder = base_folder / "parsed"
    chunks_folder = base_folder / "chunks"
    vector_folder = base_folder / "vectordb"
    memory_folder = base_folder / "memory"
    
    # Create directories if they don't exist
    try:
        parsed_folder.mkdir(parents=True, exist_ok=True)
        chunks_folder.mkdir(parents=True, exist_ok=True)
        vector_folder.mkdir(parents=True, exist_ok=True)
        memory_folder.mkdir(parents=True, exist_ok=True)
    except Exception as e:
        st.error(f"Error creating user folders: {e}")
        return None
    
    return {
        'parsed': parsed_folder,
        'chunks': chunks_folder,
        'vector': vector_folder,
        'memory': memory_folder
    }

# Get user-specific folders
USER_FOLDERS = get_user_folders(st.session_state.session_id)
if USER_FOLDERS:
    PARSED_FOLDER = USER_FOLDERS['parsed']
    CHUNKS_FOLDER = USER_FOLDERS['chunks']
else:
    st.error("Failed to create user folders")
    st.stop()

# Your existing functions with improvements
def save_markdown(content: str, filename: str, folder=PARSED_FOLDER):
    """Save markdown content to file with error handling"""
    try:
        save_path = folder / filename
        with open(save_path, "w", encoding="utf-8") as f:
            f.write(content)
        return save_path
    except Exception as e:
        st.error(f"Error saving markdown file {filename}: {e}")
        return None

def parse_pdf(file_path: Path):
    """Parse PDF file with improved error handling"""
    try:
        loader = UnstructuredPDFLoader(str(file_path))
        docs = loader.load()
        text = "\n\n".join(doc.page_content for doc in docs)
        return text.strip() if text.strip() else None
    except Exception as e:
        st.error(f"PDF parse error for {file_path.name}: {e}")
        return None

def parse_docx(file_path: Path):
    """Parse DOCX file with improved error handling"""
    api_key = os.getenv("LLAMA_CLOUD_API_KEY")
    if not api_key:
        st.error("LLAMA_CLOUD_API_KEY not found in environment variables")
        return None
    
    try:
        llama_parser = LlamaParse(
            api_key=api_key,
            result_type="markdown",
            verbose=False,  # Reduced verbosity for cleaner output
            system_prompt="Extract structured content and preserve formatting as Markdown.",
        )
        parsed = llama_parser.load_data(str(file_path))
        text = "\n\n".join(
            part.get("text", "") if isinstance(part, dict) else getattr(part, "text", getattr(part, "page_content", ""))
            for part in parsed
        )
        return text.strip() if text.strip() else None
    except Exception as e:
        st.error(f"DOCX parse error for {file_path.name}: {e}")
        return None

def parse_pptx(file_path: Path):
    """Parse PPTX file with improved error handling"""
    try:
        prs = Presentation(str(file_path))
        content = ""
        for i, slide in enumerate(prs.slides, 1):
            content += f"## Slide {i}\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content += shape.text.strip() + "\n\n"
        return content.strip() if content.strip() else None
    except Exception as e:
        st.error(f"PPTX parse error for {file_path.name}: {e}")
        return None

def parse_xlsx(file_path: Path):
    """Parse XLSX file with improved error handling"""
    return safe_read_excel(file_path)

def parse_md(file_path: Path):
    """Parse Markdown file with improved error handling"""
    return safe_read_file(file_path)

def parse_txt(file_path: Path):
    """Parse TXT file with improved error handling"""
    text = safe_read_file(file_path)
    if text:
        return f"```\n{text}\n```"
    return None

def fallback_loader(file_path: Path):
    """Fallback loader with improved error handling"""
    ext = file_path.suffix.lower()
    try:
        if ext == ".docx":
            loader = UnstructuredWordDocumentLoader(str(file_path))
        elif ext == ".pptx":
            loader = UnstructuredPowerPointLoader(str(file_path))
        elif ext == ".xlsx":
            loader = UnstructuredExcelLoader(str(file_path))
        elif ext == ".md":
            loader = UnstructuredMarkdownLoader(str(file_path))
        elif ext == ".txt":
            loader = TextLoader(str(file_path))
        else:
            return None

        docs = loader.load()
        text = "\n\n".join(doc.page_content for doc in docs)
        return text.strip() if text.strip() else None
    except Exception as e:
        st.error(f"Fallback loader error for {file_path.name}: {e}")
        return None

def parse_and_save_file(file_path: Path):
    """Parse and save file with comprehensive error handling"""
    ext = file_path.suffix.lower()
    content = None

    try:
        if ext == ".pdf":
            content = parse_pdf(file_path)
        elif ext == ".docx":
            content = parse_docx(file_path)
        elif ext == ".pptx":
            content = parse_pptx(file_path)
        elif ext == ".xlsx":
            content = parse_xlsx(file_path)
        elif ext == ".md":
            content = parse_md(file_path)
        elif ext == ".txt":
            content = parse_txt(file_path)

        if not content and ext != ".pdf":
            content = fallback_loader(file_path)

        if content:
            save_path = save_markdown(content, f"{file_path.stem}.md")
            if save_path:
                return {"page_content": content, "source": file_path.name, "save_path": save_path}
        
        st.error(f"Failed to parse {file_path.name}")
        return None
        
    except Exception as e:
        st.error(f"Error processing {file_path.name}: {e}")
        return None

def save_chunk(content: str, filename: str):
    """Save chunk content with error handling"""
    try:
        save_path = CHUNKS_FOLDER / filename
        with open(save_path, "w", encoding="utf-8") as f:
            f.write(content)
        return save_path
    except Exception as e:
        st.error(f"Error saving chunk {filename}: {e}")
        return None

def split_and_save_chunks(content: str, base_filename: str, chunk_size=1000, chunk_overlap=100):
    """Split content into chunks and save them with error handling"""
    try:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", " ", ""]
        )
        chunks = splitter.split_text(content)
        
        chunk_paths = []
        for i, chunk in enumerate(chunks, 1):
            chunk_filename = f"{base_filename}_chunk{i}.md"
            chunk_path = save_chunk(chunk, chunk_filename)
            if chunk_path:
                chunk_paths.append(chunk_path)
        
        return chunks, chunk_paths
    except Exception as e:
        st.error(f"Error splitting content for {base_filename}: {e}")
        return [], []

def initialize_vectorstore():
    """Initialize vectorstore with error handling"""
    try:
        model_name = "sentence-transformers/all-MiniLM-L6-v2"
        embeddings = HuggingFaceEmbeddings(model_name=model_name)
        
        chunk_docs = []
        for chunk_file in CHUNKS_FOLDER.glob("*.md"):
            text = safe_read_file(chunk_file)
            if text:
                chunk_docs.append(Document(page_content=text, metadata={"source": chunk_file.name}))
        
        if chunk_docs:
            vectorstore = Chroma.from_documents(
                documents=chunk_docs, 
                embedding=embeddings, 
                collection_name=f"docs_{st.session_state.session_id}",
                persist_directory=str(USER_FOLDERS['vector'])
            )
            return vectorstore
        else:
            st.warning("No valid chunks found to create vectorstore")
            return None
    except Exception as e:
        st.error(f"Error initializing vectorstore: {e}")
        return None

def setup_retrieval_chain(vectorstore):
    """Setup retrieval chain with error handling"""
    try:
        groq_api_key = os.getenv("GROQ_API_KEY")
        if not groq_api_key:
            st.error("GROQ_API_KEY not found in environment variables")
            return None
        
        llm = ChatGroq(
            api_key=groq_api_key,
            model="llama-3.1-8b-instant"
        )
        
        metadata_field_info = [
            AttributeInfo(name="source", description="Name of the document file", type="string"),
            AttributeInfo(name="file_type", description="Type of file, like pdf, docx, pptx", type="string"),
        ]
        
        document_content_description = "Content of documents which may include PDFs, Word files, slides, and text containing varied topics."
        
        retriever = SelfQueryRetriever.from_llm(
            llm=llm,
            vectorstore=vectorstore,
            document_contents=document_content_description,
            metadata_field_info=metadata_field_info,
            verbose=False,  # Reduced verbosity
        )
        
        qa_prompt = ChatPromptTemplate.from_messages([
            SystemMessagePromptTemplate.from_template(
                "You are a highly reliable assistant helping users understand documents. "
                "Always answer ONLY using the provided context and conversation history. "
                "If the answer is not present in the context, clearly state that you do not know or that the information is not available. "
                "Do NOT make up or infer facts that are not explicitly in the context."
            ),
            MessagesPlaceholder(variable_name="history"),
            HumanMessagePromptTemplate.from_template(
                """
Context:
{context}

Question:
{question}
"""
            ),
        ])
        
        output_parser = StrOutputParser()
        
        def format_docs(docs):
            if isinstance(docs, list):
                return "\n\n".join(doc.page_content for doc in docs)
            raise ValueError("Expected a list of Document objects.")
        
        retrieval_chain = RunnableSequence(
            {
                "context": RunnableLambda(lambda x: x["question"]) | retriever | RunnableLambda(format_docs),
                "question": RunnableLambda(lambda x: x["question"]),
                "history": RunnableLambda(lambda x: x.get("history", [])),
            }
            | qa_prompt
            | llm
            | output_parser
        )
        
        os.makedirs(USER_FOLDERS['memory'], exist_ok=True)
        
        def get_file_history(session_id: str) -> FileChatMessageHistory:
            file_path = USER_FOLDERS['memory'] / f"{session_id}.json"
            return FileChatMessageHistory(file_path=str(file_path))
        
        retrieval_chain_with_memory = RunnableWithMessageHistory(
            runnable=retrieval_chain,
            get_message_history=get_file_history,
            get_session_history=get_file_history,
            input_messages_key="question",
            history_messages_key="history",
        )
        
        return retrieval_chain_with_memory
    except Exception as e:
        st.error(f"Error setting up retrieval chain: {e}")
        return None

def chat_with_memory(question: str, session_id: str = "default_session"):
    """Chat with memory functionality with error handling"""
    try:
        if st.session_state.retrieval_chain_with_memory is None:
            st.error("Please upload and process documents first!")
            return None
        
        result = st.session_state.retrieval_chain_with_memory.invoke(
            {"question": question},
            config={"configurable": {"session_id": session_id}},
        )
        return result
    except Exception as e:
        st.error(f"Error during chat: {e}")
        return None

# Streamlit UI
def main():
    st.set_page_config(
        page_title="MultiDoc Memory RAG System",
        page_icon="📚",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    st.title("📚 MultiDoc Memory RAG System")
    st.markdown("---")
    
    # Sidebar for configuration
    st.sidebar.header("🔧 Configuration")
    st.sidebar.markdown(f"**Session ID:** `{st.session_state.session_id}`")
    
    # Session management
    with st.sidebar.expander("👤 Session Management"):
        if st.button("🔄 New Session"):
            # Clear current session data
            st.session_state.clear()
            st.rerun()
        
        # Show user data info
        if PARSED_FOLDER.exists():
            parsed_count = len(list(PARSED_FOLDER.glob("*.md")))
            chunks_count = len(list(CHUNKS_FOLDER.glob("*.md")))
            st.write(f"📄 Parsed docs: {parsed_count}")
            st.write(f"🔍 Chunks: {chunks_count}")
        
        if st.button("🗑️ Clear All Data"):
            try:
                user_data_path = Path("user_data") / st.session_state.session_id
                if user_data_path.exists():
                    shutil.rmtree(user_data_path)
                st.session_state.clear()
                st.success("All user data cleared!")
                st.rerun()
            except Exception as e:
                st.error(f"Error clearing data: {e}")
    
    # Check environment variables
    with st.sidebar.expander("🔑 Environment Variables"):
        llama_key = "✅" if os.getenv("LLAMA_CLOUD_API_KEY") else "❌"
        groq_key = "✅" if os.getenv("GROQ_API_KEY") else "❌"
        st.write(f"LLAMA_CLOUD_API_KEY: {llama_key}")
        st.write(f"GROQ_API_KEY: {groq_key}")
    
    # Main content area
    tab1, tab2, tab3, tab4 = st.tabs(["📁 Upload & Process", "📄 Parsed Documents", "🔍 Chunks", "💬 Chat"])
    
    with tab1:
        st.header("📁 Upload & Process Documents")
        
        uploaded_files = st.file_uploader(
            "Choose files",
            accept_multiple_files=True,
            type=['pdf', 'docx', 'pptx', 'md', 'txt', 'xlsx']
        )
        
        col1, col2 = st.columns(2)
        with col1:
            chunk_size = st.slider("Chunk Size", 500, 2000, 1000)
        with col2:
            chunk_overlap = st.slider("Chunk Overlap", 50, 500, 100)
        
        if uploaded_files:
            if st.button("🚀 Process Documents", type="primary"):
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                processed_docs = []
                total_files = len(uploaded_files)
                
                for i, uploaded_file in enumerate(uploaded_files):
                    status_text.text(f"Processing {uploaded_file.name}...")
                    
                    # Use temporary file with context manager for automatic cleanup
                    with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{uploaded_file.name}") as temp_file:
                        temp_file.write(uploaded_file.getbuffer())
                        temp_path = Path(temp_file.name)
                    
                    try:
                        # Parse the file
                        doc = parse_and_save_file(temp_path)
                        if doc:
                            processed_docs.append(doc)
                            
                            # Split into chunks
                            chunks, chunk_paths = split_and_save_chunks(
                                doc["page_content"], 
                                uploaded_file.name.rsplit(".", 1)[0],
                                chunk_size=chunk_size,
                                chunk_overlap=chunk_overlap
                            )
                            
                            if chunks:
                                st.success(f"✅ {uploaded_file.name}: {len(chunks)} chunks created")
                            else:
                                st.warning(f"⚠️ {uploaded_file.name}: No chunks created")
                    
                    except Exception as e:
                        st.error(f"Error processing {uploaded_file.name}: {e}")
                    
                    finally:
                        # Clean up temp file safely
                        safe_delete_file(temp_path)
                    
                    progress_bar.progress((i + 1) / total_files)
                
                if processed_docs:
                    status_text.text("Initializing vector store...")
                    
                    # Initialize vectorstore
                    vectorstore = initialize_vectorstore()
                    if vectorstore:
                        st.session_state.vectorstore = vectorstore
                        st.session_state.retrieval_chain_with_memory = setup_retrieval_chain(vectorstore)
                        st.success("🎉 Vector store initialized successfully!")
                    else:
                        st.error("Failed to initialize vector store")
                else:
                    st.warning("No documents were successfully processed")
                
                status_text.text("Processing complete!")
                progress_bar.progress(1.0)
    
    with tab2:
        st.header("📄 Your Parsed Documents")
        
        if PARSED_FOLDER.exists():
            parsed_files = list(PARSED_FOLDER.glob("*.md"))
            if parsed_files:
                st.info(f"📊 You have {len(parsed_files)} parsed documents in your session")
                selected_file = st.selectbox("Select a parsed document:", parsed_files)
                
                if selected_file:
                    with st.expander(f"📄 {selected_file.name}", expanded=True):
                        content = safe_read_file(selected_file)
                        if content:
                            st.markdown(content)
                            
                            # Download button
                            st.download_button(
                                label="📥 Download",
                                data=content,
                                file_name=selected_file.name,
                                mime="text/markdown"
                            )
                        else:
                            st.error("Could not read file content")
            else:
                st.info("You haven't parsed any documents yet. Please upload and process documents first.")
        else:
            st.info("You haven't parsed any documents yet. Please upload and process documents first.")
    
    with tab3:
        st.header("🔍 Your Document Chunks")
        
        if CHUNKS_FOLDER.exists():
            chunk_files = list(CHUNKS_FOLDER.glob("*.md"))
            if chunk_files:
                st.info(f"📊 You have {len(chunk_files)} chunks in your session")
                
                # Group chunks by source document
                chunk_groups = {}
                for chunk_file in chunk_files:
                    base_name = chunk_file.stem.rsplit("_chunk", 1)[0]
                    if base_name not in chunk_groups:
                        chunk_groups[base_name] = []
                    chunk_groups[base_name].append(chunk_file)
                
                # Display chunks by document
                for doc_name, chunks in chunk_groups.items():
                    with st.expander(f"📄 {doc_name} ({len(chunks)} chunks)"):
                        for chunk_file in sorted(chunks):
                            chunk_num = chunk_file.stem.split("_chunk")[-1]
                            st.subheader(f"Chunk {chunk_num}")
                            content = safe_read_file(chunk_file)
                            if content:
                                st.text_area(
                                    f"Content (Chunk {chunk_num})",
                                    content,
                                    height=150,
                                    key=f"chunk_{chunk_file.stem}_{st.session_state.session_id}"
                                )
                            else:
                                st.error(f"Could not read chunk {chunk_num}")
                            st.markdown("---")
            else:
                st.info("You haven't created any chunks yet. Please upload and process documents first.")
        else:
            st.info("You haven't created any chunks yet. Please upload and process documents first.")
    
    with tab4:
        st.header("💬 Chat with Your Documents")
        
        if st.session_state.retrieval_chain_with_memory is None:
            st.warning("Please upload and process your documents first to enable chat functionality.")
        else:
            # Chat interface
            st.subheader("Ask questions about your documents")
            
            # Display chat history
            for message in st.session_state.chat_history:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])
            
            # Chat input
            if question := st.chat_input("What would you like to know about your documents?"):
                # Add user message to chat history
                st.session_state.chat_history.append({"role": "user", "content": question})
                
                # Display user message
                with st.chat_message("user"):
                    st.markdown(question)
                
                # Get response
                with st.chat_message("assistant"):
                    with st.spinner("Thinking..."):
                        response = chat_with_memory(question, st.session_state.session_id)
                        if response:
                            st.markdown(response)
                            # Add assistant response to chat history
                            st.session_state.chat_history.append({"role": "assistant", "content": response})
                        else:
                            st.error("Sorry, I couldn't process your question. Please try again.")
            
            # Clear chat history button
            if st.button("🗑️ Clear Chat History"):
                st.session_state.chat_history = []
                st.rerun()
    
    # Footer
    st.markdown("---")
    st.markdown(
        "<div style='text-align: center; color: #666;'>"
        "MultiDoc Memory RAG System - Built with Streamlit"
        "</div>",
        unsafe_allow_html=True
    )

if __name__ == "__main__":
    main()